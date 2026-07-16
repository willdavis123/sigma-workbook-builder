#!/usr/bin/env python3
"""
use-case-slide generator
Reads a use-case-agent JSON file and produces a single-slide PPTX
matching the 10-card grid layout of the Use Case Template deck.

Usage:
    python3 generate_slide.py <input.json> <output.pptx>
"""

import sys, json, textwrap
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Sigma brand tokens ──────────────────────────────────────────────────────
SIGMA_BLUE        = RGBColor(0x1A, 0x70, 0xF1)
SIGMA_BLUE_DARK   = RGBColor(0x08, 0x36, 0x7A)
SIGMA_BLACK       = RGBColor(0x00, 0x00, 0x00)
SIGMA_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
SIGMA_GRAY_BG     = RGBColor(0xF6, 0xF6, 0xF5)
SIGMA_GRAY_RULE   = RGBColor(0xE8, 0xE6, 0xDC)
SIGMA_GRAY_MID    = RGBColor(0x76, 0x74, 0x74)
SIGMA_GRAY_MUTED  = RGBColor(0x9E, 0x9C, 0x96)
SIGMA_TEXT        = RGBColor(0x1C, 0x1E, 0x1E)

# Impact tier colors  (bg, text)
TIER_COLORS = {
    "High impact": (RGBColor(0xFA, 0xEC, 0xE7), RGBColor(0x7A, 0x2E, 0x12)),
    "Quick win":   (RGBColor(0xE6, 0xF2, 0xE0), RGBColor(0x2D, 0x60, 0x18)),
    "Insight":     (RGBColor(0xE5, 0xEE, 0xFA), RGBColor(0x0F, 0x47, 0x99)),
}

# Value tag colors  (bg, text)
TAG_COLORS = {
    "Revenue growth": (RGBColor(0xE5, 0xEE, 0xFA), RGBColor(0x08, 0x36, 0x7A)),
    "Cost avoidance": (RGBColor(0xEA, 0xF5, 0xEA), RGBColor(0x1A, 0x5C, 0x1A)),
    "Risk avoidance": (RGBColor(0xFD, 0xF0, 0xE6), RGBColor(0x7A, 0x3A, 0x00)),
    "Operations":     (RGBColor(0xF2, 0xEC, 0xFA), RGBColor(0x4A, 0x1A, 0x7A)),
    "Labor":          (RGBColor(0xFD, 0xF5, 0xE0), RGBColor(0x6B, 0x4D, 0x00)),
    "Customer":       (RGBColor(0xE0, 0xF7, 0xF5), RGBColor(0x00, 0x5F, 0x55)),
    "Insight":        (RGBColor(0xF5, 0xF5, 0xF5), RGBColor(0x3C, 0x3C, 0x3C)),
}

# ── Layout constants (inches) ───────────────────────────────────────────────
SL_W = 13.333; SL_H = 7.5          # slide dimensions
M_SIDE = 0.28                        # left/right margin
M_TOP  = 0.22                        # top margin
HDR_H  = 0.44                        # header strip height
HDR_GAP = 0.13                       # gap between header and card grid
M_BOT  = 0.18                        # bottom margin

COLS = 2; ROWS = 5
GAP_X = 0.14; GAP_Y = 0.10

GRID_TOP = M_TOP + HDR_H + HDR_GAP
GRID_W   = SL_W - 2 * M_SIDE
GRID_H   = SL_H - GRID_TOP - M_BOT
CARD_W   = (GRID_W - (COLS-1) * GAP_X) / COLS
CARD_H   = (GRID_H - (ROWS-1) * GAP_Y) / ROWS

# Card internal padding
CP = 0.07   # card padding all sides
CR = 0.06   # corner radius (used in XML)

# Internal row heights (inches)
BADGE_H    = 0.155
TITLE_H    = 0.245   # 2-line title
SUMMARY_H  = 0.205   # 2-line summary
STAT_H     = 0.150
BOTTOM_H   = 0.135
GAP_INNER  = 0.035


def inches(x): return Inches(x)
def pt(x):     return Pt(x)
def rgb_hex(c): return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def add_rounded_rect(slide, x, y, w, h, fill_rgb, corner_radius_inches=0.0, border_rgb=None, border_pt=0.5):
    """Add a filled rounded rectangle shape."""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE = 1  (rounded rect = 5)
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    # Set rounded corners via XML
    sp = shape.element
    prstGeom = sp.find('.//' + qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')
        avLst = prstGeom.find(qn('a:avLst'))
        if avLst is None:
            avLst = etree.SubElement(prstGeom, qn('a:avLst'))
        else:
            avLst.clear()
        gd = etree.SubElement(avLst, qn('a:gd'))
        # adj value: corner radius as fraction of shorter dimension * 100000
        short = min(Inches(w), Inches(h))
        adj = int(min(Inches(corner_radius_inches) / short * 100000, 50000)) if corner_radius_inches > 0 else 0
        gd.set('name', 'adj'); gd.set('fmla', f'val {adj}')

    # Fill
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb

    # Border
    line = shape.line
    if border_rgb:
        line.color.rgb = border_rgb
        line.width = Pt(border_pt)
    else:
        line.fill.background()

    shape.shadow.inherit = False
    return shape


def add_text_box(slide, x, y, w, h, text, font_name, font_size_pt, bold=False,
                  color=None, align=PP_ALIGN.LEFT, wrap=True, v_anchor=None):
    """Add a text box with a single run."""
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap

    if v_anchor:
        tf.vertical_anchor = v_anchor

    # Remove default margin
    txBox.text_frame.margin_left   = 0
    txBox.text_frame.margin_right  = 0
    txBox.text_frame.margin_top    = 0
    txBox.text_frame.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    rpr = run.font
    rpr.name = font_name
    rpr.size = Pt(font_size_pt)
    rpr.bold = bold
    if color:
        rpr.color.rgb = color
    return txBox


def add_pill_badge(slide, x, y, w, h, text, bg_rgb, text_rgb, font_size_pt=5.5):
    """Render a pill-shaped badge as rounded rect + text box."""
    # Background rounded rect
    add_rounded_rect(slide, x, y, w, h, bg_rgb, corner_radius_inches=h/2)
    # Centered text on top
    add_text_box(slide, x, y, w, h,
                 text.upper(), "DM Sans", font_size_pt,
                 bold=True, color=text_rgb, align=PP_ALIGN.CENTER)


def draw_card(slide, use_case, col, row, industry_label):
    """Draw a single use case card at grid position (col, row)."""
    # Card top-left position
    cx = M_SIDE + col * (CARD_W + GAP_X)
    cy = GRID_TOP + row * (CARD_H + GAP_Y)

    # ── Card background with border ──────────────────────────────
    add_rounded_rect(slide, cx, cy, CARD_W, CARD_H,
                     SIGMA_WHITE, corner_radius_inches=0.06,
                     border_rgb=SIGMA_GRAY_RULE, border_pt=0.75)

    # Current y cursor inside card (with padding)
    iy = cy + CP
    inner_w = CARD_W - 2 * CP
    ix = cx + CP

    # ── Row 1: impact_tier badge (left) + id number (right) ──────
    tier = use_case.get("impact_tier", "Insight")
    tier_bg, tier_fg = TIER_COLORS.get(tier, TIER_COLORS["Insight"])
    badge_w = 0.90
    add_pill_badge(slide, ix, iy, badge_w, BADGE_H,
                   tier, tier_bg, tier_fg, font_size_pt=5.5)

    # ID number — right-aligned
    id_str = str(use_case.get("id", ""))
    add_text_box(slide, ix + badge_w + 0.04, iy, inner_w - badge_w - 0.04, BADGE_H,
                 id_str, "DM Sans", 7.5, bold=True,
                 color=SIGMA_GRAY_MUTED, align=PP_ALIGN.RIGHT)

    iy += BADGE_H + GAP_INNER

    # ── Row 2: title ─────────────────────────────────────────────
    title = use_case.get("title", "")
    add_text_box(slide, ix, iy, inner_w, TITLE_H,
                 title, "DM Sans", 7.5, bold=True,
                 color=SIGMA_BLACK, wrap=True)
    iy += TITLE_H + GAP_INNER

    # ── Row 3: card_summary ───────────────────────────────────────
    summary = use_case.get("card_summary", "")
    add_text_box(slide, ix, iy, inner_w, SUMMARY_H,
                 summary, "DM Sans", 6.5, bold=False,
                 color=SIGMA_GRAY_MID, wrap=True)
    iy += SUMMARY_H + GAP_INNER

    # ── Row 4: stat box ───────────────────────────────────────────
    vs = use_case.get("value_stat", {})
    stat_num = vs.get("num", "") if isinstance(vs, dict) else str(vs)
    stat_lbl = vs.get("lbl", "") if isinstance(vs, dict) else ""

    # Light gray stat background
    add_rounded_rect(slide, ix, iy, inner_w, STAT_H,
                     SIGMA_GRAY_BG, corner_radius_inches=0.04)
    # Stat number (bold, left)
    stat_num_w = min(len(stat_num) * 0.065 + 0.1, inner_w * 0.38)
    add_text_box(slide, ix + 0.06, iy, stat_num_w, STAT_H,
                 stat_num, "DM Sans", 7.5, bold=True,
                 color=SIGMA_BLACK, align=PP_ALIGN.LEFT)
    # Stat label (muted, fills rest)
    add_text_box(slide, ix + 0.06 + stat_num_w + 0.03, iy,
                 inner_w - stat_num_w - 0.15, STAT_H,
                 stat_lbl, "DM Sans", 6.0, bold=False,
                 color=SIGMA_GRAY_MID, wrap=True)
    iy += STAT_H + GAP_INNER

    # ── Row 5: department (left) + value_tag badge (right) ────────
    dept = use_case.get("department", "")
    tag  = use_case.get("value_tag", "")
    tag_bg, tag_fg = TAG_COLORS.get(tag, (SIGMA_GRAY_BG, SIGMA_GRAY_MID))

    tag_w = min(len(tag) * 0.062 + 0.18, 1.35)
    dept_w = inner_w - tag_w - 0.06

    add_text_box(slide, ix, iy, dept_w, BOTTOM_H,
                 dept, "DM Sans", 6.0, bold=False,
                 color=SIGMA_GRAY_MUTED, wrap=False)
    add_pill_badge(slide, ix + dept_w + 0.06, iy + 0.01,
                   tag_w, BOTTOM_H - 0.02,
                   tag, tag_bg, tag_fg, font_size_pt=5.0)


def draw_header(slide, data):
    """Draw the slide header: eyebrow + company/industry title."""
    industry = data.get("industry", "")
    customer = data.get("customer", "")

    hx = M_SIDE
    hy = M_TOP
    hw = SL_W - 2 * M_SIDE

    # Eyebrow: "Example Sigma Apps in <Industry>"
    # Split long industry strings gracefully
    industry_short = industry.split(" - ")[-1] if " - " in industry else industry
    eyebrow = f"Example Sigma Apps in {industry_short}"

    # Blue left accent bar
    add_rounded_rect(slide, hx, hy + 0.04, 0.04, HDR_H - 0.10,
                     SIGMA_BLUE, corner_radius_inches=0.02)

    # Eyebrow text
    add_text_box(slide, hx + 0.10, hy, hw * 0.6, HDR_H * 0.48,
                 eyebrow, "DM Sans", 7.5, bold=False,
                 color=SIGMA_BLUE)

    # Customer name (large)
    add_text_box(slide, hx + 0.10, hy + HDR_H * 0.42, hw * 0.7, HDR_H * 0.55,
                 customer, "DM Sans", 12.5, bold=True,
                 color=SIGMA_BLACK)

    # Sigma Σ wordmark (right-aligned)
    add_text_box(slide, SL_W - M_SIDE - 1.4, hy, 1.4, HDR_H,
                 "Σ Sigma", "DM Sans", 11.0, bold=True,
                 color=SIGMA_BLUE, align=PP_ALIGN.RIGHT)


def generate_slide(json_path: str, output_path: str):
    with open(json_path) as f:
        data = json.load(f)

    use_cases = data.get("use_cases", [])
    if len(use_cases) != 10:
        print(f"Warning: expected 10 use cases, got {len(use_cases)}")

    # Create presentation
    prs = Presentation()
    prs.slide_width  = Inches(SL_W)
    prs.slide_height = Inches(SL_H)

    # Blank slide layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = SIGMA_WHITE

    # Header
    draw_header(slide, data)

    # Divider line under header
    from pptx.util import Inches as I, Pt as P
    line_shape = slide.shapes.add_connector(
        1,  # STRAIGHT = 1
        I(M_SIDE), I(M_TOP + HDR_H + 0.03),
        I(SL_W - M_SIDE), I(M_TOP + HDR_H + 0.03)
    )
    line_shape.line.color.rgb = SIGMA_GRAY_RULE
    line_shape.line.width = P(0.5)

    # Draw each card
    industry = data.get("industry", "")
    for i, uc in enumerate(use_cases[:10]):
        col = i % COLS
        row = i // COLS
        draw_card(slide, uc, col, row, industry)

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 generate_slide.py <input.json> <output.pptx>")
        sys.exit(1)
    generate_slide(sys.argv[1], sys.argv[2])

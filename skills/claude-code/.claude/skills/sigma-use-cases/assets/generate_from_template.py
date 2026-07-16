#!/usr/bin/env python3
"""
use-case-slide: template-based generator (Sigma 2026 brand template)

Populates the placeholder slide of Use_Case_Template.pptx with data from a
use-case-agent JSON file. Outputs a new single-slide PPTX with the example
slide and the Styles reference slide removed.

Usage:
    python3 generate_from_template.py <input.json> <output.pptx> [--template <template.pptx>]

The template PPTX must be assets/Use_Case_Template.pptx relative to this script,
or supplied via --template.

── Styling model (Sigma 2026 brand template) ──────────────────────────────────
The deck ships fully styled. There is exactly ONE per-card styling rule:
impact-tier color coding. For each card, the tier drives FOUR things:
  • the tier chip fill
  • the tier chip border
  • the tier chip text color
  • the left accent bar (fill + line)

Everything else (card surface, stat panel, value-tag eyebrow, titles, ids,
departments, typography) is pre-styled and IDENTICAL on every card — do NOT
recolor it. In particular, the value tag is a uniform Sigma-blue eyebrow with
no background fill; leave it exactly as shipped.

If a shape fill cannot be set for any reason, still fill the text. The card then
shows the shipped default (pale-blue chip + blue accent), which is on-brand —
only the at-a-glance tier coding is lost.
"""

import sys, os, json, shutil, argparse, re, copy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# ── Impact-tier palette (the ONE per-card rule) ─────────────────────────────
# Keys are the three canonical UPPERCASE tier strings the chip must display.
# Values give chip fill, chip border, chip text, and left accent bar colors.

TIER_STYLE = {
    "HIGH IMPACT": {
        "chip_fill":   RGBColor(0x1A, 0x70, 0xF1),
        "chip_border": RGBColor(0x1A, 0x70, 0xF1),
        "chip_text":   RGBColor(0xFF, 0xFF, 0xFF),
        "accent":      RGBColor(0x1A, 0x70, 0xF1),
    },
    "QUICK WIN": {
        "chip_fill":   RGBColor(0xEE, 0xF3, 0xFF),
        "chip_border": RGBColor(0xDC, 0xE6, 0xFE),
        "chip_text":   RGBColor(0x16, 0x60, 0xD6),
        "accent":      RGBColor(0x7F, 0xA8, 0xF0),
    },
    "INSIGHT": {
        "chip_fill":   RGBColor(0xFF, 0xFF, 0xFF),
        "chip_border": RGBColor(0xD6, 0xD6, 0xDB),
        "chip_text":   RGBColor(0x5F, 0x5F, 0x66),
        "accent":      RGBColor(0xC9, 0xC9, 0xD0),
    },
}

# Default tier when the JSON value can't be matched/normalized.
_DEFAULT_TIER = "INSIGHT"


def normalize_tier(raw: str) -> str:
    """
    Map an arbitrary impact_tier string to one of the three canonical
    UPPERCASE chip values. Tolerant of case/spacing variants.
    """
    if not raw:
        return _DEFAULT_TIER
    s = re.sub(r"\s+", " ", raw.strip()).upper()
    if s in TIER_STYLE:
        return s
    # Heuristic mapping: strategic bets → HIGH IMPACT; fast/cheap → QUICK WIN;
    # analytical/visibility plays → INSIGHT.
    if "HIGH" in s or "IMPACT" in s or "STRATEGIC" in s:
        return "HIGH IMPACT"
    if "QUICK" in s or "WIN" in s or "FAST" in s:
        return "QUICK WIN"
    if "INSIGHT" in s or "VISIBILITY" in s or "ANALY" in s:
        return "INSIGHT"
    return _DEFAULT_TIER


# ── Per-card child shape indices (within each card group) ───────────────────
# [0]  card surface  (white #FFFFFF, 1px #E5E5E9 border)  — leave as is
# [1]  left accent bar          ← tier-colored (fill + line)
# [2]  tier chip background      ← tier-colored (fill + line)
# [3]  tier chip text            ← text replace + tier text color
# [4]  id text                   ← text replace
# [5]  title text                ← text replace
# [6]  card_summary text         ← text replace
# [7]  stat panel bg (#EEF3FF)   — leave as is
# [8]  value_stat_num text       ← text replace
# [9]  value_stat_lbl text       ← text replace
# [10] department text           ← text replace
# [11] secondary panel bg        — leave as is
# [12] value_tag text (eyebrow)  ← text replace ONLY (no fill, uniform blue)

IDX_ACCENT    = 1
IDX_CHIP_BG   = 2
IDX_TIER_TEXT = 3
IDX_ID        = 4
IDX_TITLE     = 5
IDX_SUMMARY   = 6
IDX_STAT_NUM  = 8
IDX_STAT_LBL  = 9
IDX_DEPT      = 10
IDX_TAG_TEXT  = 12


# ── Helpers ─────────────────────────────────────────────────────────────────

def set_shape_fill(shape, rgb: RGBColor):
    """Set a shape's solid fill color. Best-effort; never raises."""
    try:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = rgb
    except Exception as e:
        print(f"  (warn) could not set fill: {e}", file=sys.stderr)


def set_shape_line(shape, rgb: RGBColor):
    """Set a shape's line (border) color. Best-effort; never raises."""
    try:
        line = shape.line
        line.color.rgb = rgb
    except Exception as e:
        print(f"  (warn) could not set line: {e}", file=sys.stderr)


def set_text_color(shape, rgb: RGBColor):
    """Set all run font colors in a shape's text frame."""
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = rgb


def get_full_text(shape) -> str:
    """Return the full concatenated text of all runs across all paragraphs."""
    if not shape.has_text_frame:
        return ""
    return ''.join(
        run.text
        for para in shape.text_frame.paragraphs
        for run in para.runs
    )


def set_full_text(shape, new_text: str):
    """
    Replace a shape's text, preserving the formatting of the first run/paragraph
    and clearing everything else. Enables normAutofit so text shrinks to fit
    rather than overflowing.
    """
    tf = shape.text_frame
    first_para = tf.paragraphs[0]

    first_run_rpr = None
    if first_para.runs:
        first_run_rpr = first_para.runs[0]._r.find(qn('a:rPr'))

    first_pPr = first_para._p.find(qn('a:pPr'))
    txBody = tf._txBody

    # Enable normAutofit (shrink text to fit) on the text body properties
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        for tag in ['a:normAutofit', 'a:noAutofit', 'a:spAutoFit']:
            el = bodyPr.find(qn(tag))
            if el is not None:
                bodyPr.remove(el)
        etree.SubElement(bodyPr, qn('a:normAutofit'))

    # Remove all existing <a:p> elements
    for p_el in txBody.findall(qn('a:p')):
        txBody.remove(p_el)

    # Build a single new paragraph with a single run
    new_p = etree.SubElement(txBody, qn('a:p'))
    if first_pPr is not None:
        new_p.append(copy.deepcopy(first_pPr))

    new_r = etree.SubElement(new_p, qn('a:r'))
    if first_run_rpr is not None:
        new_r.append(copy.deepcopy(first_run_rpr))

    new_t = etree.SubElement(new_r, qn('a:t'))
    new_t.text = new_text

    if first_run_rpr is not None:
        end_rpr = copy.deepcopy(first_run_rpr)
        end_rpr.tag = qn('a:endParaRPr')
        new_p.append(end_rpr)


def get_card_groups(slide):
    """Return the card groups (top-level GROUP shapes) on the slide."""
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.GROUP]


def populate_card(group, use_case: dict):
    """Fill one card group with data and apply the tier color rule."""
    children = list(group.shapes)
    if len(children) < 13:
        print(f"  (warn) card group has {len(children)} children, expected 13; "
              f"styling indices may be off", file=sys.stderr)

    uc = use_case
    tier_key = normalize_tier(uc.get("impact_tier", ""))
    style = TIER_STYLE[tier_key]

    vs = uc.get("value_stat", {})
    stat_num = vs.get("num", "") if isinstance(vs, dict) else str(vs)
    stat_lbl = vs.get("lbl", "") if isinstance(vs, dict) else ""

    # ── The one styling rule: tier color coding ──
    # Left accent bar (fill + line)
    set_shape_fill(children[IDX_ACCENT], style["accent"])
    set_shape_line(children[IDX_ACCENT], style["accent"])
    # Tier chip background (fill + border)
    set_shape_fill(children[IDX_CHIP_BG], style["chip_fill"])
    set_shape_line(children[IDX_CHIP_BG], style["chip_border"])
    # Tier chip text (content + color)
    set_full_text(children[IDX_TIER_TEXT], tier_key)
    set_text_color(children[IDX_TIER_TEXT], style["chip_text"])

    # ── Text-only fills (no recoloring — pre-styled, uniform on every card) ──
    set_full_text(children[IDX_ID], str(uc.get("id", "")))
    set_full_text(children[IDX_TITLE], uc.get("title", ""))
    set_full_text(children[IDX_SUMMARY], uc.get("card_summary", ""))
    set_full_text(children[IDX_STAT_NUM], stat_num)
    set_full_text(children[IDX_STAT_LBL], stat_lbl)
    set_full_text(children[IDX_DEPT], uc.get("department", ""))
    # Value tag — uniform blue eyebrow, NO fill change. Text only.
    set_full_text(children[IDX_TAG_TEXT], str(uc.get("value_tag", "")).upper())


def populate_industry_title(slide, data: dict):
    """Replace {industry} in the slide title shape."""
    industry = data.get("industry", "")
    industry_short = industry.split(" - ")[-1] if " - " in industry else industry
    title_text = f"Example Sigma Apps in {industry_short}"

    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.GROUP and shape.has_text_frame:
            full = get_full_text(shape)
            if "{industry}" in full:
                set_full_text(shape, title_text)
                return
    # Fallback: match a non-group shape that mentions "industry"
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.GROUP and shape.has_text_frame:
            if "industry" in get_full_text(shape).lower():
                set_full_text(shape, title_text)
                return


def build_slide_notes(slide, data: dict):
    """Write structured seller talking points into the slide notes pane."""
    use_cases = sorted(data.get("use_cases", []), key=lambda u: u["id"])
    customer = data.get("customer", "")
    lines = []

    lines.append(f"SELLER NOTES — {customer}")
    lines.append("=" * 60)
    lines.append(
        "Use these talking points if a prospect asks for more detail "
        "on any card. Each section maps to a numbered card on the slide."
    )
    lines.append("")

    for uc in use_cases:
        vs = uc.get("value_stat", {})
        stat_n = vs.get("num", "") if isinstance(vs, dict) else str(vs)
        stat_l = vs.get("lbl", "") if isinstance(vs, dict) else ""
        wb = uc.get("writeback_model", [])
        comps = uc.get("sigma_components", [])
        formula = uc.get("value_case", {}).get("directional_formula", "")
        c_start = uc.get("maturity_start", "L0")
        c_end = uc.get("maturity_target", "L2")
        cx = uc.get("complexity", "")

        lines.append(f"{uc['id']}  {uc['title'].upper()}")
        lines.append(f"    {uc['impact_tier']}  ·  {uc['value_tag']}  ·  {uc['department']}")
        lines.append("")
        lines.append(f"  The problem")
        lines.append(f"    {uc.get('pain_point', '')}")
        lines.append("")
        lines.append(f"  What Sigma does")
        lines.append(f"    {uc.get('solution', '')}")
        lines.append("")
        lines.append(f"  What gets written back")
        for w in wb:
            lines.append(f"    • {w['name']} — {w.get('description', '')}")
        lines.append("")
        lines.append(f"  Why it matters financially")
        lines.append(f"    {formula}")
        lines.append("")
        lines.append(f"  Benchmark")
        lines.append(f"    {stat_n}  {stat_l}")
        lines.append("")
        lines.append(f"  Build details")
        lines.append(f"    Complexity {cx}/5  ·  {c_start} → {c_end}  ·  {', '.join(comps)}")
        lines.append("")
        lines.append("-" * 60)
        lines.append("")

    notes_text = "\n".join(lines).rstrip()
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


def generate_from_template(json_path: str, output_path: str, template_path: str):
    with open(json_path) as f:
        data = json.load(f)

    use_cases = data.get("use_cases", [])
    if len(use_cases) != 10:
        print(f"Warning: expected 10 use cases, got {len(use_cases)}", file=sys.stderr)

    shutil.copy2(template_path, output_path)
    prs = Presentation(output_path)

    # Locate the placeholder slide: the one containing the {industry} token.
    template_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.GROUP and shape.has_text_frame:
                if "{industry}" in get_full_text(shape):
                    template_idx = i
                    break
        if template_idx is not None:
            break
    if template_idx is None:
        print("Error: could not find the placeholder slide ({industry} token)",
              file=sys.stderr)
        sys.exit(1)

    slide = prs.slides[template_idx]

    populate_industry_title(slide, data)

    groups = get_card_groups(slide)
    uc_by_id = {uc["id"]: uc for uc in use_cases}

    for group in groups:
        id_shape = list(group.shapes)[IDX_ID]
        raw_id_text = get_full_text(id_shape)
        m = re.search(r'\{(\d+)_', raw_id_text)
        if not m:
            print(f"Warning: could not parse card id from '{raw_id_text}'", file=sys.stderr)
            continue
        card_num = int(m.group(1))
        uc = uc_by_id.get(card_num)
        if uc is None:
            print(f"Warning: no use case with id={card_num}", file=sys.stderr)
            continue
        populate_card(group, uc)

    build_slide_notes(slide, data)

    # Keep ONLY the populated placeholder slide; remove all others
    # (example + Styles reference).
    prs_xml = prs.element
    sldIdLst = prs_xml.find(qn('p:sldIdLst'))
    slide_els = sldIdLst.findall(qn('p:sldId'))
    for i, el in enumerate(list(slide_els)):
        if i != template_idx:
            sldIdLst.remove(el)

    prs.save(output_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("json_path")
    parser.add_argument("output_path")
    parser.add_argument("--template", default=None,
                        help="Path to template PPTX (default: assets/Use_Case_Template.pptx beside this script)")
    args = parser.parse_args()

    if args.template:
        template_path = args.template
    else:
        script_dir = os.path.dirname(os.path.abspath(__file__))
        template_path = os.path.join(script_dir, "Use_Case_Template.pptx")

    if not os.path.exists(template_path):
        print(f"Error: template not found at {template_path}", file=sys.stderr)
        sys.exit(1)

    generate_from_template(args.json_path, args.output_path, template_path)
